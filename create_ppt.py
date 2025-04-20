from pptx import Presentation
from pptx.util import Pt

slides = [
    {
        "title": "프로젝트 개요",
        "content": [
            "• Six Thinking Hats 사고기법 기반 AI 분석 서비스 개발",
            "• 입력 → AI 분석 → 결과 웹 UI/이메일/PDF 제공",
            "• 목표: 실전 업무에서 AI Agent(windsurf) 활용성 검증, 개발 자동화 및 업무 효율화"
        ]
    },
    {
        "title": "개발 환경 및 구조",
        "content": [
            "• Frontend: React, npm",
            "• Backend: Python (FastAPI 등), REST API",
            "• 배포 예정: Vercel(프론트), PythonAnywhere(백엔드)",
            "• 폴더 구조: frontend(UI), backend(API)"
        ]
    },
    {
        "title": "AI Agent(windsurf) 활용 방식",
        "content": [
            "• AI Agent(windsurf)가 개발 전체 자동화",
            "  - 폴더/파일 구조 설계, 코드 생성, UI/UX, API/PDF/이메일 기능",
            "• 사람의 역할: 요구사항 설명, 결과 확인 및 테스트"
        ]
    },
    {
        "title": "주요 기능 및 구현 결과",
        "content": [
            "• 웹 UI: 입력창, 결과 출력, 결과 복사, PDF 첨부 이메일 전송",
            "• API: 자연어 입력 → AI 분석 → 결과 반환, PDF 변환 및 첨부",
            "• 이메일: 분석 결과 및 PDF 첨부 자동 전송",
            "• 테스트: 모든 기능 로컬에서 정상 동작 확인"
        ]
    },
    {
        "title": "업무 자동화/AI Agent의 가치",
        "content": [
            "• 개발 속도: 이틀 만에 MVP 완성",
            "• 품질: 실제 회사 테스트에서 '매우 만족' 평가",
            "• 효율성: 반복적/루틴 업무 자동화, 개발자 리소스 절감",
            "• 확장성: 다양한 사내 프로젝트에 적용 가능"
        ]
    },
    {
        "title": "시사점 및 향후 계획",
        "content": [
            "• AI Agent의 실전 활용성 검증",
            "• 배포 및 사용자 피드백 수집 예정",
            "• 보안/성능 추가 점검 및 고도화",
            "• 사내 다양한 프로젝트로 확장 검토"
        ]
    },
    {
        "title": "결론",
        "content": [
            "• windsurf AI agent는 실질적 업무 자동화, 개발 혁신, 빠른 프로토타입 제작, 높은 만족도를 실현",
            "• 앞으로 더 많은 사내 프로젝트에서 적극 활용 가치 있음"
        ]
    },
    {
        "title": "[부록] 주요 화면 캡처",
        "content": [
            "• (여기에 실제 서비스 화면 캡처 이미지를 삽입하세요)"
        ]
    }
]

prs = Presentation()

for slide in slides:
    slide_layout = prs.slide_layouts[1]
    sld = prs.slides.add_slide(slide_layout)
    sld.shapes.title.text = slide["title"]
    content = "\n".join(slide["content"])
    sld.placeholders[1].text = content
    for paragraph in sld.placeholders[1].text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(20)

prs.save("Six_Thinking_Hats_AI_발표자료.pptx")
print("PPT 파일이 생성되었습니다: Six_Thinking_Hats_AI_발표자료.pptx")